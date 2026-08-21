import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Clean Professional Color Palette
    C_NAVY_DARK = RGBColor(15, 23, 42)      # #0F172A (Deep Slate/Navy for title)
    C_NAVY_HEADER = RGBColor(30, 41, 59)    # #1E293B (Header Banner / Dark Container)
    C_WHITE = RGBColor(255, 255, 255)       # #FFFFFF
    C_BG_LIGHT = RGBColor(248, 250, 252)    # #F8FAFC (Clean Light Background)
    C_CARD_BG = RGBColor(255, 255, 255)     # #FFFFFF (White Card Fill)
    C_CARD_BORDER = RGBColor(226, 232, 240) # #E2E8F0 (Subtle Card Border)
    C_TEXT_MAIN = RGBColor(30, 41, 59)      # #1E293B (Primary Dark Text)
    C_TEXT_MUTED = RGBColor(100, 116, 139)  # #64748B (Secondary Muted Text)
    C_BLUE_ACCENT = RGBColor(37, 99, 235)   # #2563EB (Royal Blue Accent)
    C_SKY_ACCENT = RGBColor(2, 132, 199)    # #0284C7 (Sky Blue Accent)
    C_GREEN_ACCENT = RGBColor(5, 150, 105)  # #059669 (Emerald Green Accent)
    C_RED_ACCENT = RGBColor(220, 38, 38)    # #DC2626 (Crimson Red Accent)

    def set_slide_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_slide_header(slide, slide_num_str, title_text, category_text="COURSE PROJECT — REVIEW I"):
        # Header background banner
        header_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        header_rect.fill.solid()
        header_rect.fill.fore_color.rgb = C_NAVY_HEADER
        header_rect.line.fill.background()

        # Accent bar at bottom of header
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.06), Inches(13.333), Inches(0.04))
        bar.fill.solid()
        bar.fill.fore_color.rgb = C_BLUE_ACCENT
        bar.line.fill.background()

        # Category / Tracker text
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(10.0), Inches(0.3))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = C_SKY_ACCENT
        p.font.name = "Arial"

        # Title text
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.42), Inches(10.0), Inches(0.55))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = C_WHITE
        p2.font.name = "Arial"

        # Slide Number Badge
        num_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.5), Inches(0.25), Inches(1.1), Inches(0.55))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = C_BLUE_ACCENT
        num_box.line.fill.background()
        tf_num = num_box.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = f"Slide {slide_num_str}"
        p_num.alignment = PP_ALIGN.CENTER
        p_num.font.size = Pt(13)
        p_num.font.bold = True
        p_num.font.color.rgb = C_WHITE
        p_num.font.name = "Arial"

    def add_card(slide, left, top, width, height, bg_color=C_CARD_BG, border_color=C_CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1)
        else:
            card.line.fill.background()
        return card

    # =========================================================================
    # SLIDE 1: Title & Group Members' Names
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1, C_NAVY_DARK)

    # Decorative top bar
    top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = C_BLUE_ACCENT
    top_bar.line.fill.background()

    # Category pill
    cat_pill = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(0.75), Inches(4.5), Inches(0.42))
    cat_pill.fill.solid()
    cat_pill.fill.fore_color.rgb = C_NAVY_HEADER
    cat_pill.line.color.rgb = C_SKY_ACCENT
    cat_pill.line.width = Pt(1)
    p_cat = cat_pill.text_frame.paragraphs[0]
    p_cat.text = "CONTROL SYSTEMS ENGINEERING — COURSE PROJECT"
    p_cat.alignment = PP_ALIGN.CENTER
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = C_SKY_ACCENT
    p_cat.font.name = "Arial"

    # Main Title Box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.45), Inches(11.333), Inches(1.8))
    tf_t = title_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = "Cascade (Dual-Loop) Voltage-Current\nControl of a Buck Converter"
    p_t.font.size = Pt(34)
    p_t.font.bold = True
    p_t.font.color.rgb = C_WHITE
    p_t.font.name = "Arial"

    p_sub = tf_t.add_paragraph()
    p_sub.text = "Hardware, Simulation, Firmware & Real-Time Control Verification"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = C_TEXT_MUTED
    p_sub.font.name = "Arial"
    p_sub.space_before = Pt(10)

    # Presentation Review Banner Box
    rev_card = add_card(slide1, Inches(1.0), Inches(3.6), Inches(11.333), Inches(0.75), bg_color=C_NAVY_HEADER, border_color=C_BLUE_ACCENT)
    tf_rev = rev_card.text_frame
    p_rev = tf_rev.paragraphs[0]
    p_rev.text = "PROJECT REVIEW I  |  Date: 19.08.2026 (Wednesday)  |  Time: 3.45 PM Onwards"
    p_rev.alignment = PP_ALIGN.CENTER
    p_rev.font.size = Pt(14)
    p_rev.font.bold = True
    p_rev.font.color.rgb = C_WHITE
    p_rev.font.name = "Arial"

    # Metadata Cards (Group Members & Details)
    # Card 1: Group Members
    g_card = add_card(slide1, Inches(1.0), Inches(4.6), Inches(5.5), Inches(2.3), bg_color=C_NAVY_HEADER, border_color=C_CARD_BORDER)
    tf_g = g_card.text_frame
    tf_g.margin_left = Inches(0.3)
    tf_g.margin_top = Inches(0.2)
    p_g_head = tf_g.paragraphs[0]
    p_g_head.text = "GROUP MEMBERS"
    p_g_head.font.size = Pt(12)
    p_g_head.font.bold = True
    p_g_head.font.color.rgb = C_SKY_ACCENT
    p_g_head.font.name = "Arial"

    members = [
        "1. Student Member 1 (Roll No: XXXXXX01)",
        "2. Student Member 2 (Roll No: XXXXXX02)",
        "3. Student Member 3 (Roll No: XXXXXX03)",
        "4. Student Member 4 (Roll No: XXXXXX04)"
    ]
    for m in members:
        pm = tf_g.add_paragraph()
        pm.text = m
        pm.font.size = Pt(12.5)
        pm.font.color.rgb = C_WHITE
        pm.font.name = "Arial"
        pm.space_before = Pt(4)

    # Card 2: Course & Department Details
    d_card = add_card(slide1, Inches(6.833), Inches(4.6), Inches(5.5), Inches(2.3), bg_color=C_NAVY_HEADER, border_color=C_CARD_BORDER)
    tf_d = d_card.text_frame
    tf_d.margin_left = Inches(0.3)
    tf_d.margin_top = Inches(0.2)
    p_d_head = tf_d.paragraphs[0]
    p_d_head.text = "PROJECT DETAILS"
    p_d_head.font.size = Pt(12)
    p_d_head.font.bold = True
    p_d_head.font.color.rgb = C_SKY_ACCENT
    p_d_head.font.name = "Arial"

    details = [
        ("Course:", "Control Systems Engineering"),
        ("Project Code:", "CSE-PRJ-2026"),
        ("Department:", "Electrical & Electronics Engineering"),
        ("Faculty Guide:", "Project Supervisor / Evaluator")
    ]
    for label, val in details:
        pd = tf_d.add_paragraph()
        pd.text = f"{label} "
        pd.font.size = Pt(12.5)
        pd.font.bold = True
        pd.font.color.rgb = C_SKY_ACCENT
        pd.font.name = "Arial"
        pd.space_before = Pt(4)
        
        run = pd.add_run()
        run.text = val
        run.font.bold = False
        run.font.color.rgb = C_WHITE

    # =========================================================================
    # SLIDE 2: Contents
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2, C_BG_LIGHT)
    add_slide_header(slide2, "2", "Table of Contents", "OVERVIEW OF PRESENTATION")

    agenda_items = [
        ("01", "Introduction & Background", "Context of DC-DC buck converters & motivation for cascade control topology."),
        ("02", "Literature Review", "Comparative analysis of VMC, PCMC, and Digital ACMC techniques."),
        ("03", "Gap Analysis", "Identification of limitations in standard control methods & proposed solution bridge."),
        ("04", "Project Objectives", "Core engineering goals: transfer functions, safety clamping, firmware & simulation."),
        ("05", "Block Diagram & Architecture", "Cascade dual-loop topology, inner/outer loops, transfer functions & parameters."),
        ("06", "Additional Work: Hardware & Firmware", "Circuit schematic, Arduino Timer 1 Fast PWM, 10 kHz ISR & ANSI C control module."),
        ("07", "Additional Work: Simulation & Verification", "Python ODE numerical plant simulator, transient step load & overcurrent test results."),
        ("08", "References", "Key academic literature, textbooks, and application reports.")
    ]

    card_w = Inches(5.6)
    card_h = Inches(1.2)
    top_pos = [Inches(1.4), Inches(2.8), Inches(4.2), Inches(5.6)]

    for idx, (num, title, desc) in enumerate(agenda_items):
        col = idx // 4  # 0 or 1
        row = idx % 4   # 0 to 3
        left = Inches(0.8) if col == 0 else Inches(6.933)
        top = top_pos[row]

        card = add_card(slide2, left, top, card_w, card_h)

        # Number circle badge
        badge = slide2.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.2), top + Inches(0.25), Inches(0.7), Inches(0.7))
        badge.fill.solid()
        badge.fill.fore_color.rgb = C_BLUE_ACCENT
        badge.line.fill.background()
        p_b = badge.text_frame.paragraphs[0]
        p_b.text = num
        p_b.alignment = PP_ALIGN.CENTER
        p_b.font.size = Pt(14)
        p_b.font.bold = True
        p_b.font.color.rgb = C_WHITE
        p_b.font.name = "Arial"

        # Text box
        tbox = slide2.shapes.add_textbox(left + Inches(1.0), top + Inches(0.15), card_w - Inches(1.1), card_h - Inches(0.3))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT_MAIN
        p1.font.name = "Arial"

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.font.name = "Arial"
        p2.space_before = Pt(3)

    # =========================================================================
    # SLIDE 3: Introduction
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3, C_BG_LIGHT)
    add_slide_header(slide3, "3", "Introduction & Background", "DC-DC BUCK CONVERTER CONTROL")

    # Left Panel: Buck Converter Fundamentals
    panel1 = add_card(slide3, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.5))
    tf1 = panel1.text_frame
    tf1.margin_left = Inches(0.3)
    tf1.margin_top = Inches(0.3)
    tf1.margin_right = Inches(0.3)
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "DC-DC BUCK CONVERTER BASICS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_BLUE_ACCENT
    p.font.name = "Arial"

    bullet_pts_1 = [
        ("Purpose & Operation:", "Steps down a higher DC input voltage (12V) to a stable lower output voltage (5V) using Pulse Width Modulation (PWM) duty cycle manipulation."),
        ("Energy Storage Stage:", "Uses an Inductor (L = 100 µH) for current smoothing and an Output Capacitor (C = 470 µF) to filter voltage ripple."),
        ("Application Scope:", "Ubiquitous in computer power supplies, battery chargers, robotics, and automotive electronic control units (ECUs)."),
        ("Single-Loop Drawbacks:", "Traditional single-loop controllers measure only output voltage Vout. Due to output capacitor delays, they react slowly to load surges and offer ZERO short-circuit current protection.")
    ]

    for title, desc in bullet_pts_1:
        pb = tf1.add_paragraph()
        pb.text = f"• {title} "
        pb.font.size = Pt(12)
        pb.font.bold = True
        pb.font.color.rgb = C_TEXT_MAIN
        pb.font.name = "Arial"
        pb.space_before = Pt(10)
        r = pb.add_run()
        r.text = desc
        r.font.bold = False
        r.font.color.rgb = C_TEXT_MUTED

    # Right Panel: Why Dual-Loop Cascade Control?
    panel2 = add_card(slide3, Inches(6.933), Inches(1.4), Inches(5.6), Inches(5.5))
    tf2 = panel2.text_frame
    tf2.margin_left = Inches(0.3)
    tf2.margin_top = Inches(0.3)
    tf2.margin_right = Inches(0.3)
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "CASCADE DUAL-LOOP CONTROL SOLUTION"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_BLUE_ACCENT
    p.font.name = "Arial"

    bullet_pts_2 = [
        ("Nested Control Hierarchy:", "Splits control action into two nested feedback loops running at different execution frequencies."),
        ("Outer Voltage Loop (Slow - 1 kHz / 1 ms):", "Monitors Vout against reference Vref (5.0V). Computes required inductor current target (Iref). Output is HARD CLAMPED at 1.5A in software."),
        ("Inner Current Loop (Fast - 10 kHz / 100 µs):", "Monitors instantaneous inductor current IL, compares with Iref, and adjusts PWM duty cycle directly."),
        ("Timescale Separation Benefit:", "Making the inner current loop 10× faster than the outer voltage loop simplifies LC filter dynamics into an easily controlled 1st-order system for the outer loop.")
    ]

    for title, desc in bullet_pts_2:
        pb = tf2.add_paragraph()
        pb.text = f"• {title} "
        pb.font.size = Pt(12)
        pb.font.bold = True
        pb.font.color.rgb = C_TEXT_MAIN
        pb.font.name = "Arial"
        pb.space_before = Pt(10)
        r = pb.add_run()
        r.text = desc
        r.font.bold = False
        r.font.color.rgb = C_TEXT_MUTED

    # =========================================================================
    # SLIDE 4: Literature Review
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4, C_BG_LIGHT)
    add_slide_header(slide4, "4", "Literature Review", "COMPARATIVE CONTROL ANALYSIS")

    lit_cards = [
        ("Single-Loop Voltage-Mode Control (VMC)", "TRADITIONAL TOPOLOGY", C_RED_ACCENT, [
            ("Operating Mechanism", "Directly senses Vout and compares with Vref to generate PWM duty cycle via a single voltage error amplifier."),
            ("System Dynamics", "Governed by second-order LC resonant double pole: f_0 = 1 / (2π√(LC)) ≈ 734 Hz."),
            ("Key Limitations", "Severe phase lag (180°) limits loop bandwidth; vulnerable to input voltage fluctuations; NO overcurrent protection during load short circuits.")
        ]),
        ("Peak Current-Mode Control (PCMC)", "ANALOG HYBRID TOPOLOGY", C_SKY_ACCENT, [
            ("Operating Mechanism", "Outer voltage loop sets peak current threshold; inner loop turns off PWM switch when peak inductor current hits threshold."),
            ("System Dynamics", "Eliminates inductor dynamics from outer loop, reducing system to 1st-order behavior."),
            ("Key Limitations", "Highly sensitive to switching noise; suffers from subharmonic oscillation for D > 50%, requiring artificial slope compensation ramps.")
        ]),
        ("Digital Average Current-Mode Control (ACMC)", "PROPOSED CASCADE TOPOLOGY", C_GREEN_ACCENT, [
            ("Operating Mechanism", "Discrete-time dual PI loops in microcontroller; inner loop averages inductor current over switching period."),
            ("System Dynamics", "Strict bandwidth decoupling (f_ci = 2000 Hz inner vs f_cv = 200 Hz outer). Inner loop absorbs disturbances."),
            ("Key Advantages", "Superior noise immunity; direct software current limit clamping (1.5A); feedforward duty cycle calculation eliminates input line transients.")
        ])
    ]

    card_w = Inches(3.644)
    card_h = Inches(5.5)
    lefts = [Inches(0.8), Inches(4.844), Inches(8.888)]

    for idx, (title, tag, accent, points) in enumerate(lit_cards):
        left = lefts[idx]
        card = add_card(slide4, left, Inches(1.4), card_w, card_h)

        stripe = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.4), card_w, Inches(0.1))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent
        stripe.line.fill.background()

        tf = card.text_frame
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.word_wrap = True

        p_tag = tf.paragraphs[0]
        p_tag.text = tag
        p_tag.font.size = Pt(9)
        p_tag.font.bold = True
        p_tag.font.color.rgb = accent
        p_tag.font.name = "Arial"

        p_head = tf.add_paragraph()
        p_head.text = title
        p_head.font.size = Pt(13)
        p_head.font.bold = True
        p_head.font.color.rgb = C_TEXT_MAIN
        p_head.font.name = "Arial"
        p_head.space_before = Pt(4)

        for p_title, p_desc in points:
            pb = tf.add_paragraph()
            pb.text = f"• {p_title}: "
            pb.font.size = Pt(11)
            pb.font.bold = True
            pb.font.color.rgb = C_TEXT_MAIN
            pb.font.name = "Arial"
            pb.space_before = Pt(10)
            r = pb.add_run()
            r.text = p_desc
            r.font.bold = False
            r.font.color.rgb = C_TEXT_MUTED

    # =========================================================================
    # SLIDE 5: Gap Analysis
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5, C_BG_LIGHT)
    add_slide_header(slide5, "5", "Gap Analysis & Proposed Solutions", "IDENTIFYING & BRIDGING SYSTEM LIMITATIONS")

    gaps = [
        ("GAP 1: Inadequate Short-Circuit & Overcurrent Protection",
         "Existing Voltage-Mode Control (VMC) relies exclusively on output voltage feedback. During a short circuit, inductor current ramps up to destructive levels before Vout drops enough for the controller to react, leading to MOSFET thermal destruction.",
         "PROPOSED SOLUTION: Software Hard Current Clamping",
         "The outer voltage PI loop generates an intermediate reference current command (Iref) that is hard clamped at 1.5A in microcontroller firmware. Inductor current physically cannot exceed 1.5A regardless of load short circuits.",
         C_RED_ACCENT),

        ("GAP 2: LC Resonant Double-Pole Instability & Slow Recovery",
         "Single-loop control must contend with the 2nd-order LC filter transfer function G_vd(s). The 180° phase lag forces low controller gain to maintain stability, resulting in long transient recovery times (> 15 ms) during step loads.",
         "PROPOSED SOLUTION: Timescale Bandwidth Separation (10×)",
         "Designing the inner current loop with a crossover frequency (2000 Hz) ten times higher than the outer loop (200 Hz) converts the plant into a fast 1st-order current source, enabling aggressive outer loop tuning without instability.",
         C_SKY_ACCENT),

        ("GAP 3: Noise Sensitivity & Subharmonic Oscillation",
         "Peak Current-Mode Control (PCMC) suffers from noise spikes during MOSFET turn-on and subharmonic oscillations at duty cycles D > 50%, requiring complex analog slope compensation hardware.",
         "PROPOSED SOLUTION: Digital Average Control & Duty Feedforward",
         "Digital integration over 10 kHz ISR sampling averages out switching noise. Adding an explicit feedforward duty term (Vout / Vin) instantly corrects for input voltage fluctuations without waiting for error integration.",
         C_GREEN_ACCENT)
    ]

    gap_w = Inches(11.733)
    gap_h = Inches(1.7)
    gap_tops = [Inches(1.4), Inches(3.25), Inches(5.1)]

    for idx, (g_title, g_desc, s_title, s_desc, accent) in enumerate(gaps):
        top = gap_tops[idx]
        card = add_card(slide5, Inches(0.8), top, gap_w, gap_h)

        stripe = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(0.12), gap_h)
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent
        stripe.line.fill.background()

        tbox1 = slide5.shapes.add_textbox(Inches(1.1), top + Inches(0.1), Inches(5.5), gap_h - Inches(0.2))
        tf1 = tbox1.text_frame
        tf1.word_wrap = True
        tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0

        p1 = tf1.paragraphs[0]
        p1.text = g_title
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = C_RED_ACCENT if accent == C_RED_ACCENT else C_TEXT_MAIN
        p1.font.name = "Arial"

        p1_sub = tf1.add_paragraph()
        p1_sub.text = g_desc
        p1_sub.font.size = Pt(10.5)
        p1_sub.font.color.rgb = C_TEXT_MUTED
        p1_sub.font.name = "Arial"
        p1_sub.space_before = Pt(3)

        tbox2 = slide5.shapes.add_textbox(Inches(6.8), top + Inches(0.1), Inches(5.5), gap_h - Inches(0.2))
        tf2 = tbox2.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0

        p2 = tf2.paragraphs[0]
        p2.text = s_title
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = C_GREEN_ACCENT
        p2.font.name = "Arial"

        p2_sub = tf2.add_paragraph()
        p2_sub.text = s_desc
        p2_sub.font.size = Pt(10.5)
        p2_sub.font.color.rgb = C_TEXT_MUTED
        p2_sub.font.name = "Arial"
        p2_sub.space_before = Pt(3)

    # =========================================================================
    # SLIDE 6: Objectives
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide6, C_BG_LIGHT)
    add_slide_header(slide6, "6", "Project Objectives", "KEY TECHNICAL GOALS & DELIVERABLES")

    obj_cards = [
        ("01", "Mathematical Modeling & Transfer Function Derivation",
         "Derive exact system transfer functions G_id(s) = Vin/(sL) and G_vi(s) = 1/(sC). Compute PI gains for target crossover frequencies (f_ci = 2000 Hz inner, f_cv = 200 Hz outer).",
         C_BLUE_ACCENT),

        ("02", "Software Safety & Anti-Windup Integration",
         "Implement hard software current clamping at I_ref_max = 1.5A to protect hardware. Embed anti-windup clamping in both PI integrators to prevent overshoot during duty cycle saturation.",
         C_SKY_ACCENT),

        ("03", "Microcontroller Firmware & Fast PWM ISR",
         "Develop hardware-independent ANSI C control engine. Configure Arduino Timer 1 for 10 kHz Fast PWM (Mode 14, 1600 steps resolution) and 10 kHz COMPA interrupt execution.",
         C_GREEN_ACCENT),

        ("04", "Simulation Modeling & Web Studio Dashboard",
         "Build pure Python discrete-time ODE plant simulator (plant_model.py) and real-time interactive browser studio (index.html) to evaluate dynamic step load and short-circuit response.",
         C_NAVY_HEADER)
    ]

    card_w = Inches(5.6)
    card_h = Inches(2.55)
    coords = [
        (Inches(0.8), Inches(1.4)),
        (Inches(6.933), Inches(1.4)),
        (Inches(0.8), Inches(4.25)),
        (Inches(6.933), Inches(4.25))
    ]

    for idx, (num, title, desc, accent) in enumerate(obj_cards):
        left, top = coords[idx]
        card = add_card(slide6, left, top, card_w, card_h)

        stripe = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, card_w, Inches(0.08))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent
        stripe.line.fill.background()

        badge = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.3), top + Inches(0.25), Inches(0.8), Inches(0.45))
        badge.fill.solid()
        badge.fill.fore_color.rgb = accent
        badge.line.fill.background()
        pb = badge.text_frame.paragraphs[0]
        pb.text = f"OBJ {num}"
        pb.alignment = PP_ALIGN.CENTER
        pb.font.size = Pt(11)
        pb.font.bold = True
        pb.font.color.rgb = C_WHITE
        pb.font.name = "Arial"

        tbox = slide6.shapes.add_textbox(left + Inches(0.3), top + Inches(0.8), card_w - Inches(0.6), card_h - Inches(0.9))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT_MAIN
        p1.font.name = "Arial"

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.font.name = "Arial"
        p2.space_before = Pt(6)

    # =========================================================================
    # SLIDE 7: Block Diagram (Refined visual flow layout with arrows)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide7, C_BG_LIGHT)
    add_slide_header(slide7, "7", "System Block Diagram & Architecture", "CASCADE CONTROL TOPOLOGY & SIGNAL FLOW")

    # Flow Diagram Container Card
    flow_card = add_card(slide7, Inches(0.8), Inches(1.35), Inches(11.733), Inches(2.35), bg_color=C_NAVY_HEADER, border_color=C_BLUE_ACCENT)

    # Title inside card positioned cleanly at the top
    tbox_fh = slide7.shapes.add_textbox(Inches(1.0), Inches(1.45), Inches(11.333), Inches(0.35))
    tf_fh = tbox_fh.text_frame
    tf_fh.margin_left = tf_fh.margin_top = tf_fh.margin_right = tf_fh.margin_bottom = 0
    p_fh = tf_fh.paragraphs[0]
    p_fh.text = "CASCADE DUAL-LOOP FEEDBACK CONTROL FLOW DIAGRAM"
    p_fh.font.size = Pt(11)
    p_fh.font.bold = True
    p_fh.font.color.rgb = C_SKY_ACCENT
    p_fh.font.name = "Arial"

    # Block boxes and connecting arrows
    # Total width available: 11.333 inches (from 1.0 to 12.333)
    # 6 Blocks + 5 Arrows
    # Block width = 1.45 inches, Arrow width = 0.35 inches
    blocks_data = [
        ("V_ref\n(5.0V)", C_BLUE_ACCENT),
        ("Outer PI\n(1 kHz / 1 ms)", C_NAVY_DARK),
        ("Hard Clamp\n(I_ref ≤ 1.5A)", C_RED_ACCENT),
        ("Inner PI\n(10 kHz / 100µs)", C_NAVY_DARK),
        ("Duty D +\nFeedforward", C_SKY_ACCENT),
        ("Buck Hardware\n(MOSFET & LC)", C_GREEN_ACCENT)
    ]

    block_w = Inches(1.48)
    block_h = Inches(1.2)
    arrow_w = Inches(0.35)
    start_left = Inches(1.0)
    top_box = Inches(1.9)

    for i, (btext, bcolor) in enumerate(blocks_data):
        bleft = start_left + i * (block_w + arrow_w)
        
        # Block Shape
        bshape = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bleft, top_box, block_w, block_h)
        bshape.fill.solid()
        bshape.fill.fore_color.rgb = bcolor
        bshape.line.color.rgb = C_WHITE
        bshape.line.width = Pt(1)
        pt = bshape.text_frame.paragraphs[0]
        pt.text = btext
        pt.alignment = PP_ALIGN.CENTER
        pt.font.size = Pt(10)
        pt.font.bold = True
        pt.font.color.rgb = C_WHITE
        pt.font.name = "Arial"

        # Connecting Arrow (if not last)
        if i < 5:
            aleft = bleft + block_w
            arrow = slide7.shapes.add_textbox(aleft, top_box + Inches(0.35), arrow_w, Inches(0.5))
            tf_a = arrow.text_frame
            tf_a.margin_left = tf_a.margin_top = tf_a.margin_right = tf_a.margin_bottom = 0
            pa = tf_a.paragraphs[0]
            pa.text = "➔"
            pa.alignment = PP_ALIGN.CENTER
            pa.font.size = Pt(16)
            pa.font.bold = True
            pa.font.color.rgb = C_SKY_ACCENT

    # Lower Left Panel: Control Loop Specifications
    panel_l = add_card(slide7, Inches(0.8), Inches(3.9), Inches(5.6), Inches(3.1))
    tf_l = panel_l.text_frame
    tf_l.margin_left = Inches(0.3)
    tf_l.margin_top = Inches(0.2)
    tf_l.margin_right = Inches(0.3)
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "CONTROL LOOP SPECIFICATIONS & GAINS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_BLUE_ACCENT
    p.font.name = "Arial"

    loop_specs = [
        ("Inner Current Loop:", "f_sample = 10 kHz | f_c,i = 2000 Hz | Kp,i = 0.10472 | Ki,i = 22.2808"),
        ("Outer Voltage Loop:", "f_sample = 1 kHz | f_c,v = 200 Hz | Kp,v = 0.59062 | Ki,v = 125.6637"),
        ("Feedforward Term:", "Duty_ff = Vout / Vin added to duty calculation to decouple input variations."),
        ("Anti-Windup Clamping:", "Integrator accumulators restricted to prevent windup during current limit or PWM saturation.")
    ]

    for title, desc in loop_specs:
        pb = tf_l.add_paragraph()
        pb.text = f"• {title} "
        pb.font.size = Pt(11)
        pb.font.bold = True
        pb.font.color.rgb = C_TEXT_MAIN
        pb.font.name = "Arial"
        pb.space_before = Pt(5)
        r = pb.add_run()
        r.text = desc
        r.font.bold = False
        r.font.color.rgb = C_TEXT_MUTED

    # Lower Right Panel: Plant Model & Transfer Functions
    panel_r = add_card(slide7, Inches(6.933), Inches(3.9), Inches(5.6), Inches(3.1))
    tf_r = panel_r.text_frame
    tf_r.margin_left = Inches(0.3)
    tf_r.margin_top = Inches(0.2)
    tf_r.margin_right = Inches(0.3)
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "PLANT MODEL & TRANSFER FUNCTIONS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_BLUE_ACCENT
    p.font.name = "Arial"

    plant_specs = [
        ("Control-to-Current TF:", "G_id(s) = Δi_L(s) / Δd(s) ≈ V_in / (s·L)"),
        ("Current-to-Voltage TF:", "G_vi(s) = Δv_out(s) / Δi_L(s) ≈ 1 / (s·C)"),
        ("Converter Ratings:", "V_in = 12.0 V | V_ref = 5.0 V | L = 100 µH | C = 470 µF"),
        ("Safety Limits:", "Nominal Load R = 10.0 Ω | Hard Software Current Limit I_ref_max = 1.5 A")
    ]

    for title, desc in plant_specs:
        pb = tf_r.add_paragraph()
        pb.text = f"• {title} "
        pb.font.size = Pt(11)
        pb.font.bold = True
        pb.font.color.rgb = C_TEXT_MAIN
        pb.font.name = "Arial"
        pb.space_before = Pt(5)
        r = pb.add_run()
        r.text = desc
        r.font.bold = False
        r.font.color.rgb = C_TEXT_MUTED

    # =========================================================================
    # SLIDE 8: Additional Work Done (Part 1: Hardware & Firmware)
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide8, C_BG_LIGHT)
    add_slide_header(slide8, "8", "Additional Work Done: Hardware & Firmware", "HARDWARE PROTOTYPE & EMBEDDED C CODE")

    # Left Panel: Hardware Circuit Implementation
    panel_hw = add_card(slide8, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.5))
    tf_hw = panel_hw.text_frame
    tf_hw.margin_left = Inches(0.3)
    tf_hw.margin_top = Inches(0.3)
    tf_hw.margin_right = Inches(0.3)
    tf_hw.word_wrap = True

    p = tf_hw.paragraphs[0]
    p.text = "HARDWARE CIRCUIT & SENSING STAGE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_BLUE_ACCENT
    p.font.name = "Arial"

    hw_points = [
        ("Power MOSFET Switch:", "IRF540N N-Channel MOSFET driven by IR2104 / TC4420 high-speed gate driver IC to minimize switching losses."),
        ("Freewheeling Diode:", "1N5822 Schottky Barrier Rectifier providing ultra-fast recovery and low forward voltage drop (0.3V)."),
        ("LC Filter Component Selection:", "Toroidal 100 µH power inductor (high saturation current) paired with 470 µF low-ESR electrolytic capacitor."),
        ("Voltage Sensing Divider:", "Precision 10 kΩ / 4.7 kΩ divider scaling Vout (0-5V) safely to microcontroller ADC input (Pin A0)."),
        ("Current Sensing Stage:", "ACS712-05B Hall-Effect sensor (or low-side 0.1 Ω shunt resistor with differential op-amp) connected to ADC Pin A1.")
    ]

    for title, desc in hw_points:
        pb = tf_hw.add_paragraph()
        pb.text = f"• {title} "
        pb.font.size = Pt(11.5)
        pb.font.bold = True
        pb.font.color.rgb = C_TEXT_MAIN
        pb.font.name = "Arial"
        pb.space_before = Pt(8)
        r = pb.add_run()
        r.text = desc
        r.font.bold = False
        r.font.color.rgb = C_TEXT_MUTED

    # Right Panel: Microcontroller Firmware Architecture
    panel_fw = add_card(slide8, Inches(6.933), Inches(1.4), Inches(5.6), Inches(5.5))
    tf_fw = panel_fw.text_frame
    tf_fw.margin_left = Inches(0.3)
    tf_fw.margin_top = Inches(0.3)
    tf_fw.margin_right = Inches(0.3)
    tf_fw.word_wrap = True

    p = tf_fw.paragraphs[0]
    p.text = "MICROCONTROLLER FIRMWARE ARCHITECTURE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_BLUE_ACCENT
    p.font.name = "Arial"

    fw_points = [
        ("ANSI C Control Core (cascade_control.c):", "Hardware-independent modular C code structure featuring reset, update, and anti-windup functions for outer/inner loops."),
        ("Timer 1 Fast PWM Mode 14 Setup:", "Configured ICR1 = 1599 (TOP) on 16 MHz ATmega328P to generate exactly 10 kHz PWM on Pin 9 (OC1A) with 1600 steps duty resolution."),
        ("10 kHz COMPA Interrupt Service Routine:", "Executes 100 µs inner current loop inside Timer 1 ISR; divider counter triggers 1 kHz outer loop every 10th interrupt cycle."),
        ("Soft-Start Voltage Ramping:", "Firmware ramps Vref linearly over 10 ms during bootup to eliminate initial capacitor charging inrush current.")
    ]

    for title, desc in fw_points:
        pb = tf_fw.add_paragraph()
        pb.text = f"• {title} "
        pb.font.size = Pt(11.5)
        pb.font.bold = True
        pb.font.color.rgb = C_TEXT_MAIN
        pb.font.name = "Arial"
        pb.space_before = Pt(8)
        r = pb.add_run()
        r.text = desc
        r.font.bold = False
        r.font.color.rgb = C_TEXT_MUTED

    # =========================================================================
    # SLIDE 9: Additional Work Done (Part 2: Simulation & Verification)
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide9, C_BG_LIGHT)
    add_slide_header(slide9, "9", "Additional Work Done: Simulation & Verification", "NUMERICAL ODE SIMULATOR & PERFORMANCE METRICS")

    # Top Banner Box summarizing simulation tools
    sim_banner = add_card(slide9, Inches(0.8), Inches(1.35), Inches(11.733), Inches(0.85), bg_color=C_NAVY_HEADER, border_color=C_BLUE_ACCENT)
    tf_sb = sim_banner.text_frame
    tf_sb.margin_left = Inches(0.3)
    tf_sb.margin_top = Inches(0.12)
    p_sb = tf_sb.paragraphs[0]
    p_sb.text = "SIMULATION PLATFORMS & VERIFICATION ENGINES"
    p_sb.font.size = Pt(10.5)
    p_sb.font.bold = True
    p_sb.font.color.rgb = C_SKY_ACCENT
    p_sb.font.name = "Arial"

    p_sb_sub = tf_sb.add_paragraph()
    p_sb_sub.text = "1. Pure Python Discrete ODE Plant Simulator (simulation/plant_model.py)   |   2. Interactive Real-Time Web Studio (index.html)"
    p_sb_sub.font.size = Pt(12)
    p_sb_sub.font.bold = True
    p_sb_sub.font.color.rgb = C_WHITE
    p_sb_sub.font.name = "Arial"
    p_sb_sub.space_before = Pt(3)

    # 3 Key Performance Metric Cards
    metrics = [
        ("FAST TRANSIENT STARTUP", "2.68 ms", "Startup Time to 95% Vref (4.75V)", "Smooth soft-start reference ramping eliminates initial inductor/capacitor inrush current overshoot.", C_BLUE_ACCENT),
        ("PRECISE STEADY REGULATION", "5.001 V", "Steady-State Output (10Ω Nominal Load)", "Exhibits < 0.02% steady-state voltage error with negligible PWM ripple voltage.", C_GREEN_ACCENT),
        ("OVERCURRENT CLAMPING", "1.500 A", "Short-Circuit Current Clamped (1.0Ω Load)", "Under direct output short circuit, current is perfectly clamped at 1.5A safety limit without switch breakdown.", C_RED_ACCENT)
    ]

    m_w = Inches(3.644)
    m_h = Inches(4.4)
    m_lefts = [Inches(0.8), Inches(4.844), Inches(8.888)]

    for idx, (m_tag, m_val, m_title, m_desc, accent) in enumerate(metrics):
        left = m_lefts[idx]
        card = add_card(slide9, left, Inches(2.45), m_w, m_h)

        stripe = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.45), m_w, Inches(0.1))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent
        stripe.line.fill.background()

        tf = card.text_frame
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.word_wrap = True

        ptag = tf.paragraphs[0]
        ptag.text = m_tag
        ptag.font.size = Pt(9)
        ptag.font.bold = True
        ptag.font.color.rgb = accent
        ptag.font.name = "Arial"

        pval = tf.add_paragraph()
        pval.text = m_val
        pval.font.size = Pt(32)
        pval.font.bold = True
        pval.font.color.rgb = C_TEXT_MAIN
        pval.font.name = "Arial"
        pval.space_before = Pt(8)

        ptitle = tf.add_paragraph()
        ptitle.text = m_title
        ptitle.font.size = Pt(12)
        ptitle.font.bold = True
        ptitle.font.color.rgb = C_TEXT_MAIN
        ptitle.font.name = "Arial"
        ptitle.space_before = Pt(4)

        pdesc = tf.add_paragraph()
        pdesc.text = m_desc
        pdesc.font.size = Pt(11)
        pdesc.font.color.rgb = C_TEXT_MUTED
        pdesc.font.name = "Arial"
        pdesc.space_before = Pt(8)

    # =========================================================================
    # SLIDE 10: References
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide10, C_BG_LIGHT)
    add_slide_header(slide10, "10", "References & Bibliography", "ACADEMIC LITERATURE & TECHNICAL DOCUMENTATION")

    refs = [
        ("1", "Erickson, R. W., & Maksimovic, D. (2001).", "Fundamentals of Power Electronics (2nd ed.). Springer Science & Business Media.", "Foundational theory on switch-mode power converter transfer functions and average current-mode control modeling."),
        ("2", "Mohan, N., Undeland, T. M., & Robbins, W. P. (2003).", "Power Electronics: Converters, Applications, and Design (3rd ed.). John Wiley & Sons.", "Comprehensive principles of DC-DC buck converter topology design, magnetic component sizing, and semiconductor gate drive circuits."),
        ("3", "Corradini, L., Maksimović, D., Mattavelli, P., & Zane, R. (2015).", "Digital Control of High-Frequency Switched-Mode Power Supplies. IEEE Press / John Wiley & Sons.", "Discrete-time digital controller design, anti-windup integration strategies, and bandwidth separation guidelines for dual-loop converters."),
        ("4", "Texas Instruments. (2015).", "Understanding and Applying Current-Mode Control Theory. Application Report SLUA110.", "Detailed frequency response comparisons between voltage-mode and current-mode control techniques."),
        ("5", "Microchip Technology. (2020).", "Digital Control of DC/DC Converters Using Microcontrollers. Application Note AN1475.", "Implementation guides for Timer PWM peripherals, interrupt-driven ISR loops, and ADC current sensing calibration.")
    ]

    ref_w = Inches(11.733)
    ref_h = Inches(1.0)
    ref_tops = [Inches(1.4), Inches(2.55), Inches(3.7), Inches(4.85), Inches(6.0)]

    for idx, (num, authors, book_title, annotation) in enumerate(refs):
        top = ref_tops[idx]
        card = add_card(slide10, Inches(0.8), top, ref_w, ref_h)

        # Number circle badge
        badge = slide10.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), top + Inches(0.2), Inches(0.6), Inches(0.6))
        badge.fill.solid()
        badge.fill.fore_color.rgb = C_BLUE_ACCENT
        badge.line.fill.background()
        p_b = badge.text_frame.paragraphs[0]
        p_b.text = num
        p_b.alignment = PP_ALIGN.CENTER
        p_b.font.size = Pt(12)
        p_b.font.bold = True
        p_b.font.color.rgb = C_WHITE
        p_b.font.name = "Arial"

        # Content Text box
        tbox = slide10.shapes.add_textbox(Inches(1.8), top + Inches(0.12), ref_w - Inches(1.2), ref_h - Inches(0.24))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p1 = tf.paragraphs[0]
        p1.text = f"{authors} "
        p1.font.size = Pt(11.5)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT_MAIN
        p1.font.name = "Arial"

        r = p1.add_run()
        r.text = book_title
        r.font.bold = False
        r.font.italic = True
        r.font.color.rgb = C_BLUE_ACCENT

        p2 = tf.add_paragraph()
        p2.text = f"Key relevance: {annotation}"
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.font.name = "Arial"
        p2.space_before = Pt(2)

    # Save presentation
    output_filename = "Cascade_Buck_Converter_Review1_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation successfully saved to '{output_filename}'.")

if __name__ == "__main__":
    create_presentation()
